from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Define Colors
    DARK_BG = RGBColor(26, 26, 26)      # #1A1A1A (Dark Grey)
    NEON_YELLOW = RGBColor(255, 225, 53) # #FFE135 (Nano Banana Yellow)
    WHITE_TEXT = RGBColor(255, 255, 255)
    
    slides_content = [
        {
            "layout": "Title",
            "title": "Nano Banana\n(나노바나나)",
            "subtitle": "미래를 바꿀 초소형 과일 혁명\nSmall Change, Big Impact",
            "content": ""
        },
        {
            "layout": "Content",
            "title": "나노바나나란 무엇인가?",
            "content": [
                "• 나노 기술(Nanotechnology)과 생명공학의 결합",
                "• 기존 바나나의 맛과 영양을 유지하면서 크기는 1/1000로 축소",
                "• 언제 어디서나 간편하게 섭취 가능한 미래형 식량",
                "• '작지만 확실한 행복'의 결정체"
            ]
        },
        {
            "layout": "Content",
            "title": "Feature 1: Nano Size, Huge Taste",
            "content": [
                "• 쌀알 크기의 바나나 한 개에 일반 바나나 10개 분량의 에너지 압축",
                "• 분자 요리 기법을 응용한 텍스처 보존",
                "• 주머니 속에 1년치 식량 휴대 가능",
                "• 탄소 발자국 99% 감소 효과로 친환경적"
            ]
        },
        {
            "layout": "Content",
            "title": "Feature 2: Smart Peeling Tech",
            "content": [
                "• 껍질을 까는 불편함 제거: 자동 분해되는 바이오 폴리머 껍질",
                "• 공기 접촉 시 3초 만에 껍질이 기화되어 사라짐",
                "• 쓰레기 발생 Zero 실현",
                "• 위생적이고 안전한 보관 가능"
            ]
        },
        {
            "layout": "Comparison",
            "title": "일반 바나나 vs 나노바나나",
            "content_left": "Regular Banana\n\n• 무게: 120g\n• 보관: 1주일\n• 섭취: 1분\n• 휴대성: 불편함",
            "content_right": "Nano Banana\n\n• 무게: 0.1g\n• 보관: 10년 (상온)\n• 섭취: 0.1초\n• 휴대성: 무한대"
        },
        {
            "layout": "Content",
            "title": "활용 분야: 우주에서 사막까지",
            "content": [
                "• 우주비행사(Space Snacks): 부피 최소화, 영양 극대화",
                "• 재난 구호 물품: 드론으로 수백만 개 투하 가능",
                "• 스포츠 에너지원: 마라톤, 철인3종 경기 중 즉시 섭취",
                "• 다이어트: 뇌를 속여 포만감을 주는 신호 전달물질 함유"
            ]
        },
        {
            "layout": "Content",
            "title": "Quantum Potassium (양자 칼륨)",
            "content": [
                "• 기존 칼륨보다 체내 흡수율 500% 증가",
                "• 나노 캡슐화된 비타민 B6, C, 마그네슘",
                "• 즉각적인 근육 회복 및 피로 개선 효과",
                "• 섭취 시 행복 호르몬(세로토닌) 부스터 효과"
            ]
        },
        {
            "layout": "Content",
            "title": "Market Analysis",
            "content": [
                "• 글로벌 슈퍼푸드 시장의 새로운 패러다임 제시",
                "• 2030년 예상 시장 규모: 500억 달러",
                "• 타겟: 바쁜 현대인, 운동선수, 캠핑족, 우주 매니아",
                "• 경쟁사: 마이크로 멜론, 피코 파인애플 (개발 중)"
            ]
        },
        {
            "layout": "Content",
            "title": "Future Roadmap",
            "content": [
                "• 2026 Q3: 시제품 출시 및 FDA 승인 신청",
                "• 2027 Q1: 글로벌 편의점 입점 (계산대 옆 캡슐 형태)",
                "• 2028: 맛 커스텀 서비스 (초코, 딸기 등 다양한 맛)",
                "• 2030: 가정용 나노바나나 재배기 보급 시작"
            ]
        },
        {
            "layout": "Title",
            "title": "결론: 바나나의 재정의",
            "subtitle": "나노바나나는 단순한 과일이 아닙니다.\n인류의 식생활을 바꿀 거대한 혁신입니다.\n\nSmall Change, Big Impact.\n지금 나노바나나 혁명에 동참하세요!",
            "content": ""
        }
    ]

    for i, slide_data in enumerate(slides_content):
        # Create blank slide for custom styling
        slide_layout = prs.slide_layouts[6] # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # 1. Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG
        
        # 2. Design Elements (Shapes)
        # Top-left accent bar
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.5), Inches(0.2), Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NEON_YELLOW
        shape.line.fill.background() # No outline
        
        # Bottom-right geometric circle (decoration)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(6.5), Inches(2), Inches(2))
        circle.fill.transparent = True
        circle.line.color.rgb = NEON_YELLOW
        circle.line.width = Pt(1.5)

        # 3. Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_tf = title_shape.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = slide_data['title']
        title_p.font.name = 'Arial Black'
        title_p.font.size = Pt(36)
        title_p.font.color.rgb = NEON_YELLOW
        
        # 4. Content
        if slide_data['layout'] == 'Title':
            # Center big subtitle
            sub_shape = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
            sub_tf = sub_shape.text_frame
            sub_tf.word_wrap = True
            sub_p = sub_tf.paragraphs[0]
            sub_p.text = slide_data['subtitle']
            sub_p.font.name = 'Arial'
            sub_p.font.size = Pt(24)
            sub_p.font.color.rgb = WHITE_TEXT
            sub_p.alignment = PP_ALIGN.CENTER
            
        elif slide_data['layout'] == 'Comparison':
             # Left Box
            left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.2), Inches(5))
            ltf = left_box.text_frame
            ltf.text = slide_data['content_left']
            for p in ltf.paragraphs:
                p.font.size = Pt(18)
                p.font.color.rgb = WHITE_TEXT
            
            # Right Box (Highlighed)
            right_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(1.8), Inches(4.5), Inches(3.5))
            right_bg.fill.solid()
            right_bg.fill.fore_color.rgb = RGBColor(50, 50, 50)
            right_bg.line.color.rgb = NEON_YELLOW
            
            right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4.1), Inches(4))
            rtf = right_box.text_frame
            rtf.text = slide_data['content_right']
            for p in rtf.paragraphs:
                p.font.size = Pt(20)
                p.font.color.rgb = NEON_YELLOW
                p.font.bold = True

        else: # Regular Content
            content_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(4.5))
            tf = content_shape.text_frame
            tf.word_wrap = True
            
            for item in slide_data['content']:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(20)
                p.font.color.rgb = WHITE_TEXT
                p.space_after = Pt(14)

    # Save
    output_path = "NanoBanana_Presentation_Pro.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
